import pandas as pd
from pathlib import Path
from tqdm import tqdm
import hashlib
import logging
from wrapt_timeout_decorator import *

import pymupdf4llm
import fitz
import pypandoc
from pptx import Presentation
import openpyxl
import subprocess
import threading

from langdetect import detect_langs, LangDetectException
from pipeline.taskfactory import TaskWithInputFileMonitor

# Import zentrale Logging-Konfiguration
import sys
sys.path.insert(0, str(Path(__file__).parent.parent.parent / 'src'))
from pipeline_logging import setup_stage_logging


def convert_pdf_to_markdown(file_path, timeout_config=None):
    """Convert PDF to plain text with adaptive timeout based on file size.

    Extracts text only (no image descriptions) to avoid ONNX Runtime issues.
    Timeouts are adaptive based on file size.

    Returns: (content, page_count, extraction_method)
    - extraction_method: "pdf_text" (plain text extraction)
    """
    if timeout_config is None:
        timeout_config = {
            'pdf_timeout_per_mb': 2,
            'pdf_max_timeout_seconds': 120,
            'pdf_min_timeout_seconds': 10,
            'pdf_enable_fallback': True
        }

    file_size_mb = Path(file_path).stat().st_size / (1024 * 1024)

    timeout_seconds = max(
        timeout_config.get('pdf_min_timeout_seconds', 10),
        min(
            int(file_size_mb * timeout_config['pdf_timeout_per_mb']),
            timeout_config['pdf_max_timeout_seconds']
        )
    )

    @timeout(timeout_seconds)
    def _extract_pdf():
        doc = fitz.open(str(file_path))
        text = "\n\n".join([page.get_text() for page in doc])
        page_count = len(doc)
        doc.close()
        return text, page_count

    try:
        logging.debug(f"Converting PDF {Path(file_path).name} ({file_size_mb:.1f}MB, timeout={timeout_seconds}s)")
        text, page_count = _extract_pdf()
        logging.debug(f"Successfully extracted text from PDF with {page_count} pages")
        return text, page_count, "pdf_text"
    except TimeoutError:
        logging.error(f"PDF extraction timeout for {file_path} ({file_size_mb:.1f}MB, timeout={timeout_seconds}s)")
        raise TimeoutError(f"PDF extraction timeout: {file_path}")


def convert_docx_to_markdown(file_path, timeout_config=None):
    """Convert DOCX to markdown using pandoc with subprocess timeout (from config).

    Returns: (content, page_count, extraction_method)
    - extraction_method: "docx_pandoc"
    """
    if timeout_config is None:
        timeout_config = {'docx_timeout_seconds': 180}

    timeout_seconds = timeout_config.get('docx_timeout_seconds', 180)

    try:
        result = subprocess.run(
            ['pandoc', str(file_path), '-t', 'markdown', '-f', 'docx'],
            capture_output=True,
            text=True,
            timeout=timeout_seconds
        )
        if result.returncode != 0:
            raise RuntimeError(f"Pandoc error: {result.stderr}")
        return result.stdout, None, "docx_pandoc"
    except subprocess.TimeoutExpired:
        logging.warning(f"Pandoc timeout for {file_path} (>{timeout_seconds}s)")
        raise TimeoutError(f"Pandoc timeout for DOCX conversion: {file_path}")
    except Exception as e:
        raise


def convert_pptx_to_markdown(file_path):
    """Convert PPTX to markdown with structured slide sections.

    Returns: (content, page_count, extraction_method)
    - extraction_method: "pptx_python-pptx"
    """
    prs = Presentation(str(file_path))
    md_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        title = ""
        body_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if slide.shapes.title and shape.shape_id == slide.shapes.title.shape_id:
                    title = text
                elif text:
                    body_texts.append(text)
            if shape.has_table:
                body_texts.append(_table_to_markdown(shape.table))

        heading = f"## Slide {slide_num}: {title}" if title else f"## Slide {slide_num}"
        md_parts.append(heading)
        for t in body_texts:
            md_parts.append(t)
        md_parts.append("")

    content = "\n\n".join(md_parts)
    return content, len(prs.slides), "pptx_python-pptx"


def convert_xlsx_to_markdown(file_path):
    """Convert XLSX to markdown tables using openpyxl.

    Returns: (content, page_count, extraction_method)
    - extraction_method: "xlsx_openpyxl"
    """
    wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
    md_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        md_parts.append(f"## {sheet_name}")
        rows = list(ws.iter_rows(values_only=True))
        if rows:
            md_parts.append(_rows_to_markdown_table(rows))
        md_parts.append("")
    wb.close()
    content = "\n\n".join(md_parts)
    return content, len(wb.sheetnames), "xlsx_openpyxl"


def convert_md_to_markdown(file_path):
    """Read markdown file as-is.

    Returns: (content, page_count, extraction_method)
    - extraction_method: "md_passthrough"
    """
    with open(str(file_path), 'r', encoding='utf-8') as f:
        content = f.read()
    return content, None, "md_passthrough"


def _table_to_markdown(table):
    """Convert python-pptx Table object to markdown table."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace('|', '\\|') for cell in row.cells]
        rows.append(cells)
    return _rows_to_markdown_table(rows)


def _rows_to_markdown_table(rows):
    """Convert a list of rows to a markdown table string."""
    if not rows:
        return ""
    header = [str(c) if c is not None else "" for c in rows[0]]
    lines = ["| " + " | ".join(header) + " |"]
    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
    for row in rows[1:]:
        cells = [str(c).replace('|', '\\|') if c is not None else "" for c in row]
        while len(cells) < len(header):
            cells.append("")
        cells = cells[:len(header)]
        lines.append("| " + " | ".join(cells) + " |")
    return "\n".join(lines)


# Map file types to converter functions
converters = {
    'pdf':  convert_pdf_to_markdown,
    'pptx': convert_pptx_to_markdown,
    'md':   convert_md_to_markdown,
    'docx': convert_docx_to_markdown,
    'xlsx': convert_xlsx_to_markdown,
}


def provide_markdown_content(file_path, file_type, timeout_config=None):
    """Convert a file to markdown content. Returns (markdown_string, page_count, extraction_method).

    Each converter has its own timeout:
    - PDF: Adaptive (2s per MB, max 120s) text extraction
    - DOCX: Fixed 180s (subprocess timeout)
    - PPTX/XLSX/MD: No timeout (generally fast)

    extraction_method values:
    - pdf_text: PDF plain text extraction (via fitz)
    - docx_pandoc: DOCX via pandoc
    - pptx_python-pptx: PPTX via python-pptx
    - xlsx_openpyxl: XLSX via openpyxl
    - md_passthrough: Markdown pass-through
    """
    if timeout_config is None:
        timeout_config = {}

    converter = converters[file_type]
    file_size_mb = Path(file_path).stat().st_size / (1024 * 1024)
    logging.info(f"Converting {file_type.upper()} {Path(file_path).name} ({file_size_mb:.1f}MB)")
    try:
        # Pass timeout_config to converters that need it (PDF, DOCX)
        if file_type in ['pdf', 'docx']:
            content, page_count, extraction_method = converter(file_path, timeout_config)
        else:
            content, page_count, extraction_method = converter(file_path)
    except TimeoutError as e:
        logging.warning(f"Timeout converting {file_path}: {e}")
        content, page_count, extraction_method = None, None, None
    except Exception as e:
        logging.warning(f"Error converting {file_path}: {e}")
        content, page_count, extraction_method = None, None, None
    return content, page_count, extraction_method


class ExtractFileContent(TaskWithInputFileMonitor):
    def __init__(self, config_stage, config_global):
        super().__init__(config_stage, config_global)

        # Setup zentrale Logging-Konfiguration
        self.logger_configurator = setup_stage_logging(config_global)

        stage_param = config_stage['parameters']
        self.json_file_folder = Path(config_global['raw_data_folder'])
        self.file_name_inputs =  Path(config_global['raw_data_folder']) / stage_param['file_name_input']
        self.file_name_output =  Path(config_global['raw_data_folder']) / stage_param['file_name_output']
        self.file_folder = Path(config_global['file_folder'])
        self.content_folder = Path(config_global['content_folder'])
        self.file_types = stage_param['file_types']

        # Load timeout configuration from stage parameters
        self.timeout_config = stage_param.get('conversion_timeouts', {})

    def execute_task(self):
        # Logging wird jetzt zentral konfiguriert
        df_files = pd.read_pickle(self.file_name_inputs)

        # Filter only validated LiaScript files (if validation column exists)
        if 'pipe:is_valid_liascript' in df_files.columns:
            initial_count = len(df_files)
            df_files = df_files[df_files['pipe:is_valid_liascript'] == True]
            filtered_count = len(df_files)
            logging.info(f"Filtered files: {initial_count} -> {filtered_count} (removed {initial_count - filtered_count} invalid files)")

        if Path(self.file_name_output).exists():
            df_content = pd.read_pickle(self.file_name_output)
        else:
            df_content = pd.DataFrame()

        for file_type in self.file_types:
            if file_type not in converters:
                raise ValueError(f"Converter for file type '{file_type}' not found.")

        for index, row in tqdm(df_files.iterrows(), total=df_files.shape[0]):
            # Check if the content already exists for the file
            if df_content.shape[0] > 0:
                if df_content[df_content['pipe:ID'] == row['pipe:ID']].shape[0] > 0:
                    continue

            if row['pipe:file_type'] not in self.file_types:
                continue

            file_path = self.file_folder / (row['pipe:ID'] + "." + row['pipe:file_type'])
            try:
                content, page_count, extraction_method = provide_markdown_content(file_path, row['pipe:file_type'], self.timeout_config)
            except TimeoutError as e:
                logging.warning(f"Timeout extracting content from {file_path}: {e}")
                continue
            except (IOError, OSError) as e:
                logging.warning(f"File I/O error reading {file_path}: {e}")
                continue
            except Exception as e:
                logging.error(f"Unexpected error extracting content from {file_path}: {type(e).__name__}: {e}", exc_info=True)
                continue

            if content is None or content.strip() == "":
                continue

            # Save as .md file
            output_path = self.content_folder / (row['pipe:ID'] + ".md")
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(content)

            content_list_sample = {}
            content_list_sample['pipe:ID'] = row['pipe:ID']
            content_list_sample['pipe:file_type'] = row['pipe:file_type']
            content_list_sample['pipe:extraction_method'] = extraction_method
            hash_object = hashlib.sha256(content.encode('utf-8'))
            hex_dig = hash_object.hexdigest()
            content_list_sample['pipe:content_hash'] = hex_dig
            content_list_sample['pipe:content_pages'] = page_count
            content_list_sample['pipe:content_words'] = len(content.split())

            try:
                languages = detect_langs(content)
            except LangDetectException as e:
                logging.debug(f"Language detection failed for {row['pipe:ID']}: {e}")
                languages = []
            except Exception as e:
                logging.warning(f"Unexpected error detecting language for {row['pipe:ID']}: {type(e).__name__}: {e}")
                languages = []
            if languages:
                most_probable = max(languages, key=lambda lang: lang.prob)
                language, probability = most_probable.lang, most_probable.prob
            else:
                language, probability = None, None
            content_list_sample['pipe:most_prob_language'] = language
            content_list_sample['pipe:language_probability'] = probability

            df_aux = pd.DataFrame([content_list_sample])
            if df_aux.isna().all().all():
                print(df_aux)
                raise ValueError("Empty dataframe")

            df_content = pd.concat([ df_content, df_aux])
            df_content.to_pickle(self.file_name_output)
            # just for testing
            df_content.to_csv(self.file_name_output.with_suffix('.csv'))

        logging.info(f"Finished extracting content of {df_content.shape[0]} files")
        df_content.reset_index(drop=True, inplace=True)
        df_content.to_pickle(self.file_name_output)
